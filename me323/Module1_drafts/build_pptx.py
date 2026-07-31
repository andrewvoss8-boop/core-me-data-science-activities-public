"""Build ME 323 Module 1 lecture decks as .pptx from figure assets.

GIFs are embedded as-is so PowerPoint 365 / 2019 animate them in slideshow.
Marquee equations are pre-rendered PNGs in eq/. Inline formulas use plain text.
Run: python3 build_pptx.py
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
EMU_IN = 914400
SW, SH = 13.333, 7.5                      # 16:9 inches
INK = RGBColor(0x0b, 0x0b, 0x0b)
BLUE = RGBColor(0x2a, 0x78, 0xd6)
MUT = RGBColor(0x52, 0x51, 0x4e)
SURF = RGBColor(0xfc, 0xfc, 0xfb)
FONT = "Calibri"


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SURF


def _text(slide, x, y, w, h, runs, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, space=6.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    items = runs if isinstance(runs, list) else [runs]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        txt, lvl = (item, 0) if isinstance(item, str) else item
        p.level = lvl
        r = p.add_run()
        r.text = ("•  " + txt) if lvl == 1 else txt
        f = r.font
        f.name = FONT
        f.size = Pt(size - (2 if lvl == 1 else 0))
        f.bold = bold and lvl == 0
        f.color.rgb = color
    return tb


def _accent(slide, y=1.28, x=0.9, w=1.6):
    ln = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.06))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE
    ln.line.fill.background()
    ln.shadow.inherit = False


def _img(slide, path, max_w, max_h, cx, top):
    iw, ih = Image.open(path).size
    r = min(max_w / (iw / 96.0), max_h / (ih / 96.0), max_w / (iw / 150.0) * 1e9)
    # scale by fitting within max_w x max_h inches
    scale = min(max_w / iw, max_h / ih)
    w_in, h_in = iw * scale, ih * scale
    x = cx - w_in / 2
    pic = slide.shapes.add_picture(path, Inches(x), Inches(top), Inches(w_in), Inches(h_in))
    return pic, w_in, h_in


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    return s


def title_slide(prs, title, subtitle, tagline):
    s = blank(prs)
    bar = s.shapes.add_shape(1, Inches(0), Inches(2.55), Inches(SW), Inches(0.09))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    bar.shadow.inherit = False
    _text(s, 0.9, 2.75, 11.5, 1.9, title, 44, INK, bold=True)
    _text(s, 0.92, 4.35, 11.5, 0.7, subtitle, 22, BLUE, bold=True)
    _text(s, 0.92, 5.15, 11.5, 1.4, tagline, 18, MUT)
    return s


def content_title(s, title):
    _text(s, 0.9, 0.42, 11.6, 0.9, title, 30, INK, bold=True)
    _accent(s)


def bullets_slide(prs, title, bullets, size=20):
    s = blank(prs); content_title(s, title)
    _text(s, 0.95, 1.7, 11.4, 5.2, bullets, size, INK, space=12.0)
    return s


def image_slide(prs, title, img, caption=None, source=None, eq=None, max_h=4.75):
    s = blank(prs); content_title(s, title)
    top = 1.5
    if eq:
        _, ew, eh = _img(s, os.path.join(HERE, eq), 8.5, 0.85, SW / 2, top)
        top += eh + 0.15
    cap_h = (0.5 if caption else 0) + (0.4 if source else 0)
    avail_h = SH - top - 0.35 - cap_h
    _, w_in, h_in = _img(s, os.path.join(HERE, img), 11.6, min(max_h, avail_h), SW / 2, top)
    ytxt = top + h_in + 0.08
    if caption:
        _text(s, 0.9, ytxt, 11.6, 0.6, caption, 15, INK, align=PP_ALIGN.CENTER)
        ytxt += 0.5
    if source:
        _text(s, 0.9, ytxt, 11.6, 0.4, source, 11.5, MUT, align=PP_ALIGN.CENTER)
    return s


def eq_slide(prs, title, eq, lines):
    s = blank(prs); content_title(s, title)
    _, ew, eh = _img(s, os.path.join(HERE, eq), 9.5, 1.7, SW / 2, 2.2)
    _text(s, 1.2, 2.2 + eh + 0.5, 11.0, 2.0, lines, 19, INK, align=PP_ALIGN.CENTER, space=10.0)
    return s


def split_slide(prs, title, img, bullets, source=None, size=18):
    s = blank(prs); content_title(s, title)
    _, w_in, h_in = _img(s, os.path.join(HERE, img), 6.7, 4.9, 3.7, 1.7)
    _text(s, 7.6, 1.75, 5.3, 5.0, bullets, size, INK, space=11.0)
    if source:
        _text(s, 0.6, 1.7 + h_in + 0.05, 6.6, 0.5, source, 11, MUT, align=PP_ALIGN.CENTER)
    return s


# ---------------------------------------------------------------- Lecture 1
def build_lecture1():
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    title_slide(prs, "Designing a Beam Under Messy Physics",
                "ME 323 Module 1 — Lecture 1",
                "Pick a 3D-printed I-beam that carries the most load per gram.\n"
                "The equations give you a start, not the answer.")

    bullets_slide(prs, "The challenge", [
        ("Fixed: total height 18 mm, flange width 10 mm, test span 150 mm, printed length 172 mm, PLA.", 1),
        ("You choose: web thickness b and web height H_web (which sets flange thickness).", 1),
        ("Goal: maximize strength-to-weight.", 1),
        ("Constraint you can't dodge: each print-and-test is slow and expensive, so you get very few.", 1),
    ], size=21)
    image_slide(prs, "The challenge", "figures/fig_isection.png", max_h=4.6)

    bullets_slide(prs, "Why this is hard", [
        "Three things stand between you and a clean optimization:",
        ("Printed beams can fracture along the layer lines: the flange peels off the web at a printed interface whose strength appears in no handbook.", 1),
        ("Lateral-torsional buckling depends on the fixture, not just the beam.", 1),
        ("The failure modes overlap. Real beams fail in mixed, progressive ways.", 1),
        "We will price the layer-line mode with one calibrated number — but that number drifts with print session and support condition, and the other two problems remain. The equations narrow the problem without solving it.",
    ], size=20)

    bullets_slide(prs, "Three modeled capacity branches", [
        ("Bending: the beam yields in flexure.", 1),
        ("Flange-web separation: the shear flow at the printed junction exceeds the layer-line bond strength.", 1),
        ("Lateral-torsional buckling (LTB): the beam tips over sideways.", 1),
        "The class capacity is the plain minimum of the three. The dominant-mode label is a model proxy, not proof of the observed fracture.",
    ], size=21)

    image_slide(prs, "Bending", "figures/fig_bending_stress.png", eq="eq/bend.png", max_h=4.2)

    eq_slide(prs, "Why I-beams work", "eq/ix_integral.png", [
        "Bending strength is carried by I_x, and I_x counts each bit of area by the square of its distance from the neutral axis.",
        "Material at the edges earns its mass; material near the center barely helps.",
        "So put the mass in flanges far out, top and bottom. The web holds them apart (sets the lever arm) and carries the shear.",
        "Maximize I_x per gram: that is the strength-to-weight game this module plays.",
    ])
    image_slide(prs, "Same mass, more stiffness", "figures/fig_why_ibeams.png",
                caption="Three sections of equal area. Relocating the same material to the flanges multiplies I_x.",
                source="Even against a solid bar of the same height, the I-beam wins: its mass sits at the extreme fibers.")
    bullets_slide(prs, "The catch", [
        "The I-beam logic says: thinner web, taller section, mass pushed outward.",
        "Push too far and you wake the other two modes:",
        ("a flange-web junction too thin to carry the shear flow — the printed layer-line bond lets go,", 1),
        ("a section too tall and narrow, which tips over (LTB) before it ever yields.", 1),
        "The efficient shape and the failure modes pull against each other; resolving that tension is the design problem.",
    ], size=21)

    eq_slide(prs, "Shear in a beam: the check you already know", "eq/vqit.png", [
        "In each half-span V = P/2. Quick estimate: the web carries it all, tau_avg = V/(b*H_web). The full formula prices the longitudinal shear on ANY horizontal plane: Q(y) is the first moment of the area beyond the plane, t(y) the width there.",
        "It peaks at the neutral axis (3V/2A for a rectangle) — so that is where the textbook check compares against sigma_y/sqrt(3) = 43.9 MPa.",
        "Run it: at (b, H_web) = (2, 12) the web would not rupture until ~2495 N, but bending yields the flange at 835 N — a 3x margin.",
        "Verdict: for a beam made of ONE material, shear never governs here. Ours is not one material.",
    ])

    image_slide(prs, "Flange-web separation: shear flow at the printed junction",
                "figures/fig_shear_web.png", eq="eq/sepflow.png",
                caption="Every observed shear-type failure is a flange peeling off the web — not the most-stressed plane, the weakest one. Slide the same VQ/(I_x t) to the junction: Q becomes the flange's first moment, t the joint width — the classic built-up-beam glue-line check (V = P/2).",
                source="At the junction the stress is ~14% LOWER than at the neutral axis: this plane fails first because it is weaker, not because it is more stressed.",
                max_h=3.5)

    bullets_slide(prs, "What the separation check assumes", [
        "The stress side is textbook mechanics — the same VQ/(I_x t) you already knew, evaluated one plane up.",
        "The strength side is not: tau_i is the bond strength along the printed layer lines, weaker than the bulk plastic, listed in no handbook.",
        "Pre-lab 1 starts it at the bulk guess sigma_y/sqrt(3) = 43.9 MPa and calibrates it from the data — it lands near 18 MPa, and it drifts with print session and support condition.",
        "The check predicts a capacity trend on one candidate plane. It is a dominant-mode proxy, not a fracture diagnosis.",
    ], size=22)

    image_slide(prs, "Fracture vs yield, on the test machine", "figures/fig_brittle_traces.png",
                caption="Same fixture, same material. Thin webs shed the load in one drop; thicker webs plateau.",
                source="No equation on the previous slides distinguishes these two endings.")

    eq_slide(prs, "Lateral-torsional buckling: the idea", "eq/ltb.png", [
        "A tall, thin-flanged beam does not yield in place. It rolls sideways and twists at a load below its bending strength — a stability failure, not a strength one.",
        "The elastic critical moment, in the Eurocode C1/C2 closed form (Timoshenko elastic-stability theory).",
        "The next four slides open up every symbol, because LTB decides your optimum and rests on the most assumptions.",
        "You will not re-derive any of it: P_LTB arrives pre-coded in Pre-lab 1. The goal is to read that code knowing what each symbol assumes.",
    ])
    eq_slide(prs, "The critical moment, term by term", "eq/ltb_R.png", [
        "Under the root, three effects trade against each other:",
        "C_w / I_y  —  warping rigidity: the flanges resisting counter-bending.",
        "L_b^2 GJ / (pi^2 E I_y)  —  St. Venant torsion: the open section's twisting stiffness.",
        "(C_2 z_g)^2  —  load height: loading the top flange is destabilizing, so it subtracts.",
        "C_1 = 1.35 (central point load),  C_2 = 0.55 (load-height factor).",
    ])
    eq_slide(prs, "From critical moment to a load", "eq/ltb_cap.png", [
        "The beam cannot exceed its yield moment, so LTB is capped at yield:",
        "take the smaller of the elastic buckling moment and the yield moment,",
        "convert moment to load with M_max = PL/4, so P = 4M / L,",
        "with L_b = kL and z_g = c = TH/2 (load at the top surface).",
        "The min is a crude nod to inelastic buckling.",
    ])
    s = blank(prs); content_title(s, "What feeds it: the section constants")
    _, _, hsec = _img(s, os.path.join(HERE, "eq/sec_all.png"), 10.6, 2.9, SW / 2, 1.55)
    _text(s, 0.7, 1.55 + hsec + 0.25, 12.0, 2.6, [
        ("I_y: weak-axis stiffness — small for narrow flanges, which is why tall thin beams tip.", 1),
        ("J: St. Venant torsion, summed over web and flanges as thin rectangles.", 1),
        ("C_w: warping constant, the flanges' resistance to counter-bending.", 1),
        "All three come from section_props in Pre-lab 1.",
    ], 17, INK, space=8.0)
    bullets_slide(prs, "The assumptions underneath", [
        "The formula is exact for an idealized beam, and ours differs from it in four ways.",
        ("Thin-walled open section: J and C_w use thin-rectangle theory. Webs up to 7 mm against a 10 mm flange stretch 'thin'; J drifts.", 1),
        ("Timoshenko elastic stability: M_cr assumes linear-elastic, perfectly straight, no residual stress. A print has all three imperfections.", 1),
        ("Warping torsion: the C_w term assumes ends free to warp. The real fixture restrains them — which is what k absorbs.", 1),
        ("Elastic, isotropic material: E = 2.5 GPa, G = E/2.6 (nu = 0.3). Printed PLA is layered and anisotropic.", 1),
        "Each gap widens the error bar you should carry on the LTB capacity.",
    ], size=18)
    eq_slide(prs, "The one number that rules LTB: k", "eq/k_scale.png", [
        "k is the effective-length factor: how far the beam twists and bends sideways between whatever restrains it.",
        "Set by the fixture — how the ends grip, whether they warp, where load lands — not by the beam.",
        "Textbook k = 1 is a simply-supported fork end free to warp; stiffer restraint drops k below 1.",
        "Handbook starting value k = 0.33. Calibration on the class subset gives k = 0.377, which sets the equation design.",
        "Since M_cr ~ 1/k^2, a small change swings LTB hard — and the calibrated equation optimum sits on a bend-LTB knife edge (capacities within ~1%). k does not transfer between fixtures: confirm it on ours before betting on that optimum.",
    ])

    bullets_slide(prs, "LTB: what you must own, and what is background", [
        "You just saw Iy, J, Cw, warping, load height, and effective length. You are responsible for this much of it:",
        ("Yours: what LTB is (a tall thin web tipping sideways before the material yields), when it competes (thin b, large H_web), what k means physically, and why a calibrated k does not transfer between fixtures.", 1),
        ("Yours: reading the capacity plot — which branch governs where, and how much margin the runner-up has.", 1),
        ("Background: deriving Iy, J, Cw, or the Mcr formula. The code is supplied; you will never derive or re-implement it in this module.", 1),
        "The test of understanding: what happens to the LTB branch if the fixture braces the compression flange, and would you still trust the equation optimum?",
    ], size=19)

    image_slide(prs, "Which mode drives?", "figures/fig_failure_mode_map.png",
                caption="Tall thin webs tip (LTB). Short thin webs are separation-limited: junction shear flow beats the layer-line bond. The middle bends.")

    bullets_slide(prs, "Does the theory match the data?", [
        "Pre-lab 1 answers this with your beam dataset: compute each beam's predicted capacity and mode, then plot predictions against measurements.",
        "Watch for two things:",
        ("Points below the line: beams that broke before the equations said they would.", 1),
        ("Beams whose recorded failure looks nothing like the predicted mode.", 1),
        "[Placeholder: parity plot from the 44-beam dataset once strength_N lands.]",
    ], size=20)

    image_slide(prs, "One assumption, two different beams", "figures/fig_two_optima.png",
                caption="Illustration: handbook starting values versus the class calibration (66.8 MPa, k = 0.377, tau_i = 16.8 MPa), which moves the equation optimum to (1.10, 13.25).")

    bullets_slide(prs, "Where this leaves us", [
        "The equations are useful but partial:",
        ("they cannot tell a yield plateau from a one-drop fracture,", 1),
        ("their LTB term rests on fixture assumptions,", 1),
        ("the modes blur,", 1),
        ("and we have only a few expensive tests.", 1),
        "A data-driven model that carries its own uncertainty is built for this. Next lecture: Gaussian Processes and Bayesian optimization.",
    ], size=20)

    image_slide(prs, "The loop this module runs", "figures/fig_workflow_end2end.png",
                caption="Physics says what should happen and where failure lives. The data model says what the tests support and where we're still uncertain. Your judgment decides what to build. The test says what we got wrong.",
                max_h=4.1)

    bullets_slide(prs, "Before next time: Pre-lab 1", [
        "In ME323_Module1_Prelab1_FailureModes:",
        ("reduce four raw force-displacement traces to strength — including deciding what beam 9's \"twisted off the test stand\" peak really measures,", 1),
        ("load the beam data, compute strength-to-weight,", 1),
        ("code flexural yield and the junction shear-flow separation check (start tau_i at the bulk guess 43.9 MPa),", 1),
        ("compare dominant-mode proxies with measured strengths and failure notes, then calibrate (sigma_y, k, tau_i),", 1),
        ("optimize strength-to-weight.", 1),
        "The design you find there becomes one of the class's two ground-truth queries later — the pre-lab explains how those queries work. Bring questions.",
    ], size=20)

    out = os.path.join(HERE, "ME323_Module1_Lecture1.pptx")
    prs.save(out); return out


# ---------------------------------------------------------------- Lecture 2
def build_lecture2():
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    title_slide(prs, "From Data to Decisions Under Uncertainty",
                "ME 323 Module 1 — Lecture 2",
                "A model that learns from sparse data and tells you how much it does not know.\n"
                "Flow and examples adapted from Prof. Bilionis's data-science lecturebook.")

    bullets_slide(prs, "The equations overpredicted a strong beam", [
        ("The three calibrated capacity equations selected b = 1.10 mm, H_web = 13.25 mm and predicted 48.7 N/g.", 1),
        ("Ground-truth data from earlier tests gave 475.7 N = 37.48 N/g, about 23% below the equation prediction.", 1),
        ("All three modeled capacities tie within about 1% there, so the design sits on a mode boundary.", 1),
        "The equations found a strong beam, but overpredicted its performance. Today's tool adds a statement of confidence to the point estimate.",
    ], size=20)

    bullets_slide(prs, "You already have most of this", [
        "From ME 239:",
        ("Bayes' rule, Gaussian distribution", 1),
        ("Multivariate Gaussian and covariance", 1),
        ("Conditioning a Gaussian on data", 1),
        "New today, but only one step past 239:",
        ("Gaussian Process; acquisition functions (Bayesian optimization)", 1),
    ], size=20)

    eq_slide(prs, "Bayes, in one line", "eq/bayes.png", [
        "You start with a belief about beam strength. You test a beam. You update the belief.",
        "Prior: belief before the test. Likelihood: compatibility with the result. Posterior: belief after the test.",
        "Carry a distribution over unknown beam strength, not a single guess.",
    ])

    image_slide(prs, "A Gaussian is a belief about one beam", "figures/fig_gauss_strength_repeats.png",
                caption="Four repeat tests at one design give mean 558.2 N and fitted scatter 7.6 N, about 1.4%.",
                source="Pooled campaign repeats give about 4.4% with session/printer structure. The class uses 3% as a working assumption and checks its effect.")
    image_slide(prs, "Two beams at once: covariance", "figures/fig_mvn_correlation.png",
                caption="Nearby designs share material, geometry, and physics, so their strengths move together.",
                source="The covariance matrix writes that down. This is the piece the GP is built on.")
    image_slide(prs, "Conditioning: measure one, update the other", "figures/gif_conditioning.gif",
                caption="Measure beam A and the joint Gaussian conditions: belief about B shifts and tightens.",
                source="You did this algebra in ME 239. Everything today is this move, repeated.  (animated in slideshow)")
    image_slide(prs, "A GP links predictions across many designs", "figures/gif_mvn_to_gp.gif",
                caption="The same joint-Gaussian update can cover 2 designs or 120.",
                source="A kernel makes a test influence nearby designs more than distant ones.  (animated in slideshow)")
    image_slide(prs, "A Gaussian Process, before and after data", "figures/fig_gp_prior_posterior.png",
                caption="Before data the kernel proposes candidate curves; each test kills the ones that disagree.",
                source="The module models log(str/w): exp(mu_log) is the posterior median in N/g; the band is conditional uncertainty.")
    image_slide(prs, "Each test reshapes nearby predictions", "figures/fig_gp_current_campaign_slice.png",
                caption="Current class-campaign data on a thin-web slice: the band narrows near tested beams.",
                source="The band stays wide where the model has less support.")
    image_slide(prs, "The kernel is a modeling choice", "figures/fig_kernel_lengthscale.png",
                caption="The length scale answers: how far does one test's influence reach?",
                source="No single right answer; there are defensible and indefensible ones.")
    image_slide(prs, "Noise changes both the fit and its uncertainty", "figures/fig_noise_uncertainty_shared_scales.png",
                caption="Pre-lab 2 refits at 1%, 3%, and 10% using shared row-wise color scales.",
                source="The 1% fit recommends (1.58, 14.88); the 3% and 10% fits recommend (1.00, 13.39).")

    bullets_slide(prs, "Epistemic and aleatory uncertainty do different jobs", [
        ("Epistemic: uncertainty about the underlying average trend because tests are sparse. An informative new beam can reduce it.", 1),
        ("Aleatory: print-to-print and test-to-test scatter at one nominal design. Testing another location does not remove it.", 1),
        ("For one future observation in log space: sigma_total = sqrt(sigma_epi^2 + r^2), with r = 0.03.", 1),
        "MUI and EI use epistemic sigma because they value learnable uncertainty. A reliability bound for one future printed beam uses total sigma.",
    ], size=20)

    image_slide(prs, "One dial balances performance and uncertainty", "figures/fig_explore_exploit.png", eq="eq/mui.png",
                source="a(x): acquisition-function score; mu(x): predicted log(str/w); sigma(x): epistemic uncertainty; psi: how strongly uncertainty affects the choice.",
                max_h=3.7)

    image_slide(prs, "The loop: Bayesian optimization", "figures/gif_bo_mui.gif",
                caption="Fit, evaluate the acquisition function, choose the design with the highest score, test, and refit.",
                source="Seven tests, no gradient, no formula for the dip.  (animated in slideshow)")

    image_slide(prs, "EI weighs the chance and size of a win", "figures/fig_mui_ei_current.png",
                caption="EI asks how much a new test would beat the best tested beam so far, on average.",
                source="Increasing xi often favors less-certain regions, but does not guarantee a more exploratory recommendation.")

    bullets_slide(prs, "Zoom out: spending a test budget", [
        "Each print-and-test costs a machine slot, a test engineer hour, and days of queue. The class gets two pieces of ground-truth data from earlier tests through the frozen staff model, and each group gets one real print.",
        "Optimal experimental design is the batch version of the acquisition question: given N tests, which set teaches the most or finds the best fastest?",
        "Choosing the next test is the engineering decision.",
    ], size=21)

    image_slide(prs, "Physics-informed GP: the two tracks meet", "figures/fig_decision_map_4panel.png",
                caption="A vanilla GP sees geometry only. A physics-informed GP can also receive signals from the calibrated capacity equations.",
                source="Submission 1 compares both approaches rather than assuming either one must win.")

    bullets_slide(prs, "Decide first; run sensitivity checks afterward", [
        "You will meet six modeling \"knobs\" in Submission 1. Do not treat them as six equal choices.",
        ("Students decide: how physics enters the model and the risk posture — how predicted performance and uncertainty affect the choice.", 1),
        ("Students check afterward: kernel, noise, and target-scale sensitivity.", 1),
        "If reasonable assumptions move the recommendation, report that instability. If they point to the same neighborhood, report local robustness. An unstable result can justify using a test to reduce uncertainty.",
    ], size=19)

    image_slide(prs, "First compare what physics and the GP predict", "figures/fig_decision_map_4panel.png",
                caption="The top row compares calibrated physics and the GP posterior median on the same design axes.",
                max_h=4.5)
    image_slide(prs, "Then inspect disagreement and epistemic uncertainty", "figures/fig_decision_map_4panel.png",
                caption="The bottom row shows where the models disagree and where epistemic uncertainty remains large.",
                max_h=4.5)

    bullets_slide(prs, "The activity from here", [
        ("Pre-lab 2: fit the GP, write MUI (and EI), work the noise assumption, find the GP-recommended beam.", 1),
        ("Two class results: add ground-truth data from earlier tests for the equation design and the GP design — the same two beams for everyone.", 1),
        ("Submission 1: combine physics, GP, and the two new points; pick the beam you will print; defend it in the memo.", 1),
        ("Print, test, reflect, redesign: Submission 2 folds your own test result back into the model and asks what it changes.", 1),
    ], size=21)

    bullets_slide(prs, "Before next time: Pre-lab 2", [
        "In ME323_Module1_Prelab2_ML:",
        ("warm up by rebuilding the GP from ME 239's Normal-fit in four small steps (fit one Normal, covariance, conditioning, many beams at once),", 1),
        ("compare vanilla GP setups and read posterior median, epistemic uncertainty, and total predictive uncertainty,", 1),
        ("write both MUI and EI,", 1),
        ("refit under 1% / 3% / 10% noise and note what moves,", 1),
        ("record the locked class design: b = 1.00 mm, H_web = 13.39 mm, posterior median 36.7 N/g.", 1),
        "That beam is the class's second ground-truth data point from earlier tests. Bring your rationale, not just the number.",
    ], size=20)

    out = os.path.join(HERE, "ME323_Module1_Lecture2.pptx")
    prs.save(out); return out


if __name__ == "__main__":
    print("wrote", build_lecture1())
    print("wrote", build_lecture2())
